from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# adesso Farben
ADESSO_DARK  = RGBColor(0x0D, 0x3B, 0x6E)   # #0D3B6E
ADESSO_LIGHT = RGBColor(0x7B, 0xAF, 0xD4)   # #7BAFD4
ADESSO_ACC   = RGBColor(0x1A, 0x6F, 0xBF)   # #1A6FBF
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT   = RGBColor(0xF0, 0xF4, 0xF8)
GRAY_TEXT    = RGBColor(0x44, 0x55, 0x66)
GREEN        = RGBColor(0x22, 0xC5, 0x5E)
YELLOW       = RGBColor(0xF5, 0x9E, 0x0B)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # vollständig leer


def add_rect(slide, x, y, w, h, fill=None, line=None, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = wrap
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def bullet_box(slide, items, x, y, w, h, size=16, color=GRAY_TEXT,
               icon="▸", title=None, title_color=ADESSO_DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = title
        run.font.size = Pt(size + 2)
        run.font.bold = True
        run.font.color.rgb = title_color
        run.font.name = "Calibri"
    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"


# ─────────────────────────────────────────────
# FOLIE 1 — Titelfolie
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)

add_rect(slide, 0, 0, W, H, fill=ADESSO_DARK)
add_rect(slide, 0, H - Inches(0.6), W, Inches(0.6), fill=ADESSO_LIGHT)
add_rect(slide, Inches(0.5), Inches(2.6), Inches(0.08), Inches(2.2), fill=ADESSO_LIGHT)

add_text(slide, "adesso", Inches(0.8), Inches(2.5), Inches(6), Inches(1.2),
         size=52, bold=True, color=WHITE)
add_text(slide, "EventApp", Inches(3.55), Inches(2.5), Inches(8), Inches(1.2),
         size=52, bold=True, color=ADESSO_LIGHT)

add_text(slide, "Tinder für adesso-Events — einfach wischen, mitmachen.",
         Inches(0.8), Inches(3.75), Inches(9), Inches(0.7),
         size=20, color=RGBColor(0xBD, 0xD5, 0xEA), italic=True)

add_text(slide, "Pouegueu Fotso, Stephane  ·  Mai 2026",
         Inches(0.8), Inches(6.5), Inches(8), Inches(0.5),
         size=13, color=RGBColor(0x9B, 0xBB, 0xD6))


# ─────────────────────────────────────────────
# FOLIE 2 — Das Problem
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.35), fill=ADESSO_DARK)
add_rect(slide, 0, H - Inches(0.18), W, Inches(0.18), fill=ADESSO_LIGHT)

add_text(slide, "Das Problem", Inches(0.55), Inches(0.3), Inches(10), Inches(0.75),
         size=28, bold=True, color=WHITE)

problems = [
    ("📧", "Events werden per E-Mail oder Teams-Nachricht angekündigt — geht schnell unter"),
    ("🔍", "Kein zentraler Ort zum Entdecken von Veranstaltungen"),
    ("📍", "Standort-übergreifende Events schwer zu finden"),
    ("📋", "Kein direktes Feedback / Interesse-Signaling für Organisatoren"),
    ("🚫", "Kein Überblick über Kapazitäten und Wartelisten"),
]

for i, (icon, text) in enumerate(problems):
    y = Inches(1.6) + i * Inches(0.95)
    add_rect(slide, Inches(0.5), y, Inches(11.8), Inches(0.82),
             fill=RGBColor(0xF0, 0xF4, 0xF8))
    add_text(slide, icon, Inches(0.65), y + Pt(6), Inches(0.6), Inches(0.7), size=22)
    add_text(slide, text, Inches(1.35), y + Pt(8), Inches(10.8), Inches(0.65),
             size=16, color=GRAY_TEXT)


# ─────────────────────────────────────────────
# FOLIE 3 — Die Lösung
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.35), fill=ADESSO_DARK)
add_rect(slide, 0, H - Inches(0.18), W, Inches(0.18), fill=ADESSO_LIGHT)

add_text(slide, "Die Lösung", Inches(0.55), Inches(0.3), Inches(10), Inches(0.75),
         size=28, bold=True, color=WHITE)

add_text(slide, "adessoEventApp", Inches(0.5), Inches(1.5), Inches(12), Inches(0.75),
         size=26, bold=True, color=ADESSO_DARK, align=PP_ALIGN.CENTER)

add_text(slide, "Eine mobile App, die adesso-Mitarbeitenden hilft,\nFirmenevents zu entdecken und spontan teilzunehmen.",
         Inches(1.0), Inches(2.25), Inches(11), Inches(0.9),
         size=17, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

cards = [
    ("👆", "Wischen", "Rechts = Interesse\nLinks = Überspringen"),
    ("📍", "Standort", "Events nach\nStandort filtern"),
    ("🏆", "Typen", "Sport, Meeting,\nFreizeit u.v.m."),
    ("⏳", "Warteliste", "Automatisch bei\nvollen Events"),
]

for i, (icon, title, desc) in enumerate(cards):
    x = Inches(0.4) + i * Inches(3.15)
    add_rect(slide, x, Inches(3.3), Inches(2.9), Inches(3.5),
             fill=ADESSO_DARK)
    add_text(slide, icon, x, Inches(3.45), Inches(2.9), Inches(0.7),
             size=30, align=PP_ALIGN.CENTER)
    add_text(slide, title, x, Inches(4.2), Inches(2.9), Inches(0.55),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x, Inches(4.75), Inches(2.9), Inches(1.8),
             size=13, color=ADESSO_LIGHT, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# FOLIE 4 — Funktionen im Detail
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.35), fill=ADESSO_DARK)
add_rect(slide, 0, H - Inches(0.18), W, Inches(0.18), fill=ADESSO_LIGHT)

add_text(slide, "Funktionen im Detail", Inches(0.55), Inches(0.3), Inches(12), Inches(0.75),
         size=28, bold=True, color=WHITE)

cols = [
    ("👤  Mitarbeiter", [
        "Event-Feed mit Swipe-Geste",
        "Standort- und Typ-Filter",
        "KI-Empfehlungen (nach 5 Swipes)",
        "Eigene Teilnahmen einsehen",
        "Wartelisten-Benachrichtigung",
    ]),
    ("📁  Organisator", [
        "Events erstellen & bearbeiten",
        "Teilnehmerliste verwalten",
        "Events absagen",
        "Eigenes Dashboard",
        "Standort zuweisen",
    ]),
    ("⚙️  Admin", [
        "Alle Events moderieren",
        "Standorte verwalten",
        "Benutzer & Rollen verwalten",
        "Events aktivieren / absagen",
        "Vollständige Übersicht",
    ]),
]

for i, (title, items) in enumerate(cols):
    x = Inches(0.4) + i * Inches(4.25)
    add_rect(slide, x, Inches(1.5), Inches(4.05), Inches(5.6),
             fill=GRAY_LIGHT)
    add_rect(slide, x, Inches(1.5), Inches(4.05), Inches(0.5),
             fill=ADESSO_ACC)
    add_text(slide, title, x + Inches(0.12), Inches(1.54), Inches(3.8), Inches(0.42),
             size=14, bold=True, color=WHITE)
    bullet_box(slide, items,
               x + Inches(0.15), Inches(2.1), Inches(3.75), Inches(4.8),
               size=13, color=GRAY_TEXT, icon="✓")


# ─────────────────────────────────────────────
# FOLIE 5 — User Journey
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.35), fill=ADESSO_DARK)
add_rect(slide, 0, H - Inches(0.18), W, Inches(0.18), fill=ADESSO_LIGHT)

add_text(slide, "User Journey — So einfach geht's", Inches(0.55), Inches(0.3),
         Inches(12), Inches(0.75), size=28, bold=True, color=WHITE)

steps = [
    ("1", "Einloggen", "Mit adesso\nE-Mail anmelden"),
    ("2", "Standort\nwählen", "Beim ersten Login\nStandort auswählen"),
    ("3", "Events\nentdecken", "Feed öffnen &\nEvents durchstöbern"),
    ("4", "Wischen", "Rechts = Interesse\nLinks = Skip"),
    ("5", "Teilnehmen", "Bestätigung oder\nWarteliste"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.35) + i * Inches(2.55)
    # Kreis
    add_rect(slide, x + Inches(0.65), Inches(1.6), Inches(1.2), Inches(1.2),
             fill=ADESSO_DARK)
    add_text(slide, num, x + Inches(0.65), Inches(1.65), Inches(1.2), Inches(1.05),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x, Inches(2.9), Inches(2.45), Inches(0.7),
             size=14, bold=True, color=ADESSO_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x, Inches(3.6), Inches(2.45), Inches(0.85),
             size=12, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    # Pfeil (außer nach letztem)
    if i < 4:
        add_text(slide, "→", x + Inches(1.9), Inches(1.9), Inches(0.6), Inches(0.6),
                 size=22, bold=True, color=ADESSO_LIGHT, align=PP_ALIGN.CENTER)

# KI-Box
add_rect(slide, Inches(0.5), Inches(4.65), Inches(12.3), Inches(1.15),
         fill=RGBColor(0xEF, 0xF6, 0xFF))
add_text(slide, "🤖  KI-Personalisierung",
         Inches(0.65), Inches(4.72), Inches(4), Inches(0.45),
         size=14, bold=True, color=ADESSO_DARK)
add_text(slide, "Nach 5 Swipes analysiert die KI deine Präferenzen und sortiert den Feed automatisch nach deinen Interessen.",
         Inches(0.65), Inches(5.15), Inches(12), Inches(0.5),
         size=13, color=GRAY_TEXT)


# ─────────────────────────────────────────────
# FOLIE 6 — Technologie
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.35), fill=ADESSO_DARK)
add_rect(slide, 0, H - Inches(0.18), W, Inches(0.18), fill=ADESSO_LIGHT)

add_text(slide, "Technologie-Stack", Inches(0.55), Inches(0.3), Inches(12), Inches(0.75),
         size=28, bold=True, color=WHITE)

tech = [
    ("⚛️  Frontend", ADESSO_DARK, [
        "React 18 + TypeScript",
        "Tailwind CSS (Mobile-first)",
        "React Spring (Swipe-Animation)",
        "Vite + i18n (DE / EN)",
        "Zustand (State Management)",
    ]),
    ("🖥️  Backend", ADESSO_ACC, [
        "Node.js + Express",
        "Prisma ORM + SQLite",
        "JWT Authentication",
        "REST API",
        "KI-Scoring Service",
    ]),
    ("🚀  Deployment", RGBColor(0x05, 0x96, 0x69), [
        "Docker-ready",
        "Vite Dev-Proxy",
        "Port 3003 (Backend)",
        "Port 5173 (Frontend)",
        "Umgebungsvariablen (.env)",
    ]),
]

for i, (title, color, items) in enumerate(tech):
    x = Inches(0.4) + i * Inches(4.25)
    add_rect(slide, x, Inches(1.5), Inches(4.05), Inches(5.6), fill=GRAY_LIGHT)
    add_rect(slide, x, Inches(1.5), Inches(4.05), Inches(0.5), fill=color)
    add_text(slide, title, x + Inches(0.12), Inches(1.54), Inches(3.8), Inches(0.42),
             size=14, bold=True, color=WHITE)
    bullet_box(slide, items,
               x + Inches(0.15), Inches(2.1), Inches(3.75), Inches(4.8),
               size=13, color=GRAY_TEXT, icon="▸")


# ─────────────────────────────────────────────
# FOLIE 7 — Nächste Schritte
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.35), fill=ADESSO_DARK)
add_rect(slide, 0, H - Inches(0.18), W, Inches(0.18), fill=ADESSO_LIGHT)

add_text(slide, "Nächste Schritte & Ausblick", Inches(0.55), Inches(0.3),
         Inches(12), Inches(0.75), size=28, bold=True, color=WHITE)

roadmap = [
    ("✅", "MVP abgeschlossen", "Feed, Swipe, Rollen, Standorte, KI-Scoring, Organizer-Dashboard", GREEN),
    ("🔜", "Push-Benachrichtigungen", "Mitarbeiter erhalten Alerts bei neuen Events am eigenen Standort", YELLOW),
    ("🔜", "Kalender-Integration", "Events direkt in Outlook / Google Calendar exportieren", YELLOW),
    ("🔜", "Bewertungen & Kommentare", "Feedback nach dem Event hinterlassen", YELLOW),
    ("🔜", "Native Mobile App", "iOS / Android via React Native oder PWA", ADESSO_LIGHT),
]

for i, (icon, title, desc, color) in enumerate(roadmap):
    y = Inches(1.55) + i * Inches(1.05)
    add_rect(slide, Inches(0.5), y, Inches(11.8), Inches(0.88), fill=GRAY_LIGHT)
    add_rect(slide, Inches(0.5), y, Inches(0.08), Inches(0.88), fill=color)
    add_text(slide, icon, Inches(0.65), y + Pt(8), Inches(0.5), Inches(0.65), size=18)
    add_text(slide, title, Inches(1.25), y + Pt(6), Inches(3.5), Inches(0.4),
             size=14, bold=True, color=ADESSO_DARK)
    add_text(slide, desc, Inches(1.25), y + Pt(30), Inches(10.8), Inches(0.4),
             size=12, color=GRAY_TEXT)


# ─────────────────────────────────────────────
# FOLIE 8 — Abschluss
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=ADESSO_DARK)
add_rect(slide, 0, H - Inches(0.6), W, Inches(0.6), fill=ADESSO_LIGHT)
add_rect(slide, Inches(0.5), Inches(2.4), Inches(0.08), Inches(2.6), fill=ADESSO_LIGHT)

add_text(slide, "Danke!", Inches(0.8), Inches(2.3), Inches(11), Inches(1.3),
         size=54, bold=True, color=WHITE)
add_text(slide, "Fragen & Demo",
         Inches(0.8), Inches(3.65), Inches(11), Inches(0.7),
         size=22, color=ADESSO_LIGHT, italic=True)

add_text(slide, "adesso", Inches(0.8), Inches(5.2), Inches(3), Inches(0.65),
         size=20, bold=True, color=WHITE)
add_text(slide, "EventApp", Inches(2.18), Inches(5.2), Inches(4), Inches(0.65),
         size=20, bold=True, color=ADESSO_LIGHT)
add_text(slide, "· Pouegueu Fotso, Stephane  ·  Mai 2026",
         Inches(4.0), Inches(5.28), Inches(8), Inches(0.5),
         size=13, color=RGBColor(0x9B, 0xBB, 0xD6))


# Speichern
out = os.path.join(os.path.dirname(__file__), "adessoEventApp_Praesentation.pptx")
prs.save(out)
print(f"Gespeichert: {out}")
