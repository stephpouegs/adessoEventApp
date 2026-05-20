"""
Prototype-Präsentation mit echten App-Screenshots (2 Minuten).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Farben ───────────────────────────────────────
DARK   = RGBColor(0x0D, 0x3B, 0x6E)
LIGHT  = RGBColor(0x7B, 0xAF, 0xD4)
ACC    = RGBColor(0x1A, 0x6F, 0xBF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG     = RGBColor(0xF0, 0xF4, 0xF8)
GRAY   = RGBColor(0x44, 0x55, 0x66)
GREEN  = RGBColor(0x22, 0xC5, 0x5E)
YELLOW = RGBColor(0xF5, 0x9E, 0x0B)
PHONE_BODY = RGBColor(0x1C, 0x1C, 0x1E)

DOCS = os.path.dirname(__file__)
SS   = os.path.join(DOCS, "screenshots")

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Hilfsfunktionen ──────────────────────────────

def rect(slide, x, y, w, h, fill=None, radius=False):
    shp = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        int(x), int(y), int(w), int(h)
    )
    shp.line.fill.background()
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    return shp


def txt(slide, text, x, y, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tb.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def bullets(slide, items, x, y, w, h, size=15, color=GRAY, icon="->"):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, W, H, fill=BG)
    rect(slide, 0, 0, Inches(5.2), H, fill=DARK)
    rect(slide, Inches(5.2), H - Inches(0.18), W - Inches(5.2), Inches(0.18), fill=LIGHT)
    txt(slide, title,
        Inches(0.45), Inches(0.28), Inches(4.6), Inches(0.65),
        size=22, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.45), Inches(0.95), Inches(4.6), Inches(0.55),
            size=13, color=LIGHT, italic=True)


def phone_frame(slide, img_path, cx, cy, phone_w=Inches(2.7)):
    """Zeichnet ein Smartphone-Gehäuse und legt den Screenshot hinein."""
    ratio     = 844 / 390           # iPhone-Seitenverhältnis
    phone_h   = phone_w * ratio
    corner    = Inches(0.28)
    bar_h     = Inches(0.18)        # Statusbar oben
    chin_h    = Inches(0.14)        # Kinn unten
    border    = Inches(0.065)       # Rahmendicke

    x = cx - phone_w / 2
    y = cy - phone_h / 2

    # Gehäuse (abgerundet per Shape 5 = RoundedRectangle)
    body = slide.shapes.add_shape(5, int(x), int(y), int(phone_w), int(phone_h))
    body.adjustments[0] = 0.06
    body.fill.solid()
    body.fill.fore_color.rgb = PHONE_BODY
    body.line.fill.background()

    # Bildschirmrahmen
    scr_x = x + border
    scr_y = y + border
    scr_w = phone_w - 2 * border
    scr_h = phone_h - 2 * border

    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path,
                                 int(scr_x), int(scr_y),
                                 int(scr_w), int(scr_h))
    else:
        rect(slide, scr_x, scr_y, scr_w, scr_h,
             fill=RGBColor(0x1A, 0x1A, 0x2E))

    # Notch / Dynamic Island
    notch_w = Inches(0.55)
    notch_h = Inches(0.10)
    notch_x = x + phone_w / 2 - notch_w / 2
    notch_y = y + border + Inches(0.04)
    notch = slide.shapes.add_shape(5, int(notch_x), int(notch_y),
                                    int(notch_w), int(notch_h))
    notch.fill.solid()
    notch.fill.fore_color.rgb = PHONE_BODY
    notch.line.fill.background()

    # Seiten-Button
    btn_w = Inches(0.045)
    btn_h = Inches(0.22)
    btn = rect(slide,
               x + phone_w - border / 2 - btn_w / 2,
               y + phone_h * 0.32,
               btn_w, btn_h, fill=RGBColor(0x3A, 0x3A, 0x3C))
    return phone_w, phone_h


def talking_point(slide, items, x, y, w=Inches(5.8)):
    """Sprechpunkt-Box rechts."""
    bh = Inches(0.45) * len(items) + Inches(0.3)
    rect(slide, x, y, w, bh, fill=WHITE)
    bullets(slide, items, x + Inches(0.2), y + Inches(0.12),
            w - Inches(0.3), bh - Inches(0.2),
            size=14, color=DARK, icon="  ")


# ═══════════════════════════════════════════════
# FOLIE 1 — Titel
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, fill=DARK)
rect(slide, 0, H - Inches(0.55), W, Inches(0.55), fill=LIGHT)

# Hintergrund-Dekor
for i in range(6):
    s = Inches(2.5 + i * 0.8)
    rect(slide,
         W * 0.55 + Inches(i * 1.1), H / 2 - s / 2, s, s,
         fill=RGBColor(0x15, 0x50, 0x90))

txt(slide, "adesso", Inches(1.0), Inches(2.2), Inches(5), Inches(1.3),
    size=56, bold=True, color=WHITE)
txt(slide, "EventApp", Inches(1.0), Inches(3.4), Inches(6), Inches(1.1),
    size=56, bold=True, color=LIGHT)
txt(slide, "Events entdecken — einfach wischen.",
    Inches(1.0), Inches(4.6), Inches(7), Inches(0.65),
    size=19, color=RGBColor(0xBD, 0xD5, 0xEA), italic=True)
txt(slide, "2-Minuten Demo  |  Mai 2026",
    Inches(1.0), Inches(6.5), Inches(6), Inches(0.5),
    size=13, color=RGBColor(0x8A, 0xAA, 0xC8))

# Mini-Phones Deko
for i, fn in enumerate(["01_login.png", "02_feed.png", "03_swipe_right.png"]):
    fp = os.path.join(SS, fn)
    px = Inches(8.5) + i * Inches(1.55)
    phone_frame(slide, fp, px, H / 2, phone_w=Inches(1.3))


# ═══════════════════════════════════════════════
# FOLIE 2 — Login & Onboarding
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header(slide, "Login & Onboarding", "Schritt 1: Anmelden & Standort wählen")

phone_frame(slide, os.path.join(SS, "01_login.png"),
            Inches(2.6), H / 2 + Inches(0.3), phone_w=Inches(2.85))

txt(slide, "Einstieg",
    Inches(5.55), Inches(1.1), Inches(7.3), Inches(0.6),
    size=26, bold=True, color=DARK)

items = [
    "Anmeldung mit adesso E-Mail-Adresse",
    "Nur @adesso.de Adressen erlaubt",
    "Beim ersten Login: Standort wählen",
    "Rollen: Mitarbeiter / Organisator / Admin",
]
bullets(slide, items,
        Inches(5.55), Inches(1.85), Inches(7.0), Inches(3.5),
        size=17, color=GRAY, icon="  •")

rect(slide, Inches(5.55), Inches(5.0), Inches(7.3), Inches(0.75),
     fill=DARK)
txt(slide, "  Sicher & einfach — kein Passwort nötig",
    Inches(5.55), Inches(5.05), Inches(7.2), Inches(0.65),
    size=14, bold=True, color=LIGHT)


# ═══════════════════════════════════════════════
# FOLIE 3 — Event-Feed
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header(slide, "Event-Feed entdecken", "Schritt 2: Events browsen & filtern")

phone_frame(slide, os.path.join(SS, "02_feed.png"),
            Inches(2.6), H / 2 + Inches(0.3), phone_w=Inches(2.85))

txt(slide, "Der Feed",
    Inches(5.55), Inches(1.1), Inches(7.3), Inches(0.6),
    size=26, bold=True, color=DARK)

features = [
    ("📍", "Standort-Filter", "Nur Events in meiner Stadt"),
    ("🏷️", "Typ-Chips",       "Sport, Meeting, Freizeit …"),
    ("🤖", "KI-Ranking",      "Personalisiert nach Swipe-Verlauf"),
    ("📅", "Datum & Kapazität","Direkt auf der Karte sichtbar"),
]
for i, (icon, title, desc) in enumerate(features):
    y = Inches(1.8) + i * Inches(1.15)
    rect(slide, Inches(5.55), y, Inches(7.3), Inches(1.0), fill=BG)
    txt(slide, icon, Inches(5.7), y + Inches(0.22), Inches(0.5), Inches(0.55), size=20)
    txt(slide, title, Inches(6.3), y + Inches(0.1), Inches(3.5), Inches(0.42),
        size=15, bold=True, color=DARK)
    txt(slide, desc, Inches(6.3), y + Inches(0.52), Inches(6.3), Inches(0.42),
        size=13, color=GRAY)


# ═══════════════════════════════════════════════
# FOLIE 4 — Swipe-Aktion
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header(slide, "Wischen & Entscheiden", "Schritt 3: Rechts = Dabei, Links = Skip")

phone_frame(slide, os.path.join(SS, "03_swipe_right.png"),
            Inches(2.6), H / 2 + Inches(0.3), phone_w=Inches(2.85))

txt(slide, "So einfach geht Teilnahme",
    Inches(5.55), Inches(1.1), Inches(7.5), Inches(0.7),
    size=24, bold=True, color=DARK)

# Rechts-Swipe Box
rect(slide, Inches(5.55), Inches(1.95), Inches(3.4), Inches(1.3),
     fill=RGBColor(0xDC, 0xFC, 0xE8))
txt(slide, "Rechts wischen",
    Inches(5.75), Inches(2.05), Inches(3.0), Inches(0.45),
    size=15, bold=True, color=GREEN)
txt(slide, "Ich nehme teil\n(oder Warteliste bei vollen Events)",
    Inches(5.75), Inches(2.5), Inches(3.0), Inches(0.65),
    size=13, color=GRAY)

# Links-Swipe Box
rect(slide, Inches(9.15), Inches(1.95), Inches(3.4), Inches(1.3),
     fill=RGBColor(0xFE, 0xE2, 0xE2))
txt(slide, "Links wischen",
    Inches(9.35), Inches(2.05), Inches(3.0), Inches(0.45),
    size=15, bold=True, color=RGBColor(0xDC, 0x26, 0x26))
txt(slide, "Kein Interesse\n(Event erscheint nicht wieder)",
    Inches(9.35), Inches(2.5), Inches(3.0), Inches(0.65),
    size=13, color=GRAY)

# Buttons
for i, (label, color) in enumerate([
    ("⟲  Reload", RGBColor(0xF5, 0x9E, 0x0B)),
    ("✕  Ablehnen", RGBColor(0xDC, 0x26, 0x26)),
    ("♥  Dabei", GREEN),
    ("📍  Map", ACC),
]):
    bx = Inches(5.55) + i * Inches(1.88)
    rect(slide, bx, Inches(3.6), Inches(1.75), Inches(0.72), fill=DARK)
    txt(slide, label, bx, Inches(3.65), Inches(1.75), Inches(0.62),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(slide, "Oder Button-Leiste unten nutzen:",
    Inches(5.55), Inches(3.4), Inches(7.5), Inches(0.38),
    size=12, color=GRAY, italic=True)

# Warteliste-Hinweis
rect(slide, Inches(5.55), Inches(4.55), Inches(7.3), Inches(0.85),
     fill=RGBColor(0xFF, 0xF9, 0xC4))
txt(slide, "  Warteliste: Wenn ein Event voll ist, kommt man automatisch auf die Warteliste und wird informiert.",
    Inches(5.55), Inches(4.6), Inches(7.2), Inches(0.75),
    size=13, color=RGBColor(0x78, 0x35, 0x00))


# ═══════════════════════════════════════════════
# FOLIE 5 — Meine Events & Profil
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header(slide, "Meine Events & Profil", "Eigene Teilnahmen und Einstellungen")

# Zwei Phones nebeneinander
phone_frame(slide, os.path.join(SS, "04_my_events.png"),
            Inches(1.65), H / 2 + Inches(0.3), phone_w=Inches(2.45))
phone_frame(slide, os.path.join(SS, "05_profile.png"),
            Inches(4.25), H / 2 + Inches(0.3), phone_w=Inches(2.45))

# Labels unter Phones
txt(slide, "Meine Events", Inches(0.5), Inches(6.6), Inches(2.4), Inches(0.4),
    size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
txt(slide, "Profil", Inches(3.1), Inches(6.6), Inches(2.4), Inches(0.4),
    size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)

txt(slide, "Persönlicher Bereich",
    Inches(6.1), Inches(1.1), Inches(6.8), Inches(0.6),
    size=24, bold=True, color=DARK)

items_my = [
    "Alle bestätigten Teilnahmen auf einen Blick",
    "Tap auf Event -> vollständige Details & Abmelden",
    "Eigene Events (für Organisatoren)",
]
items_prof = [
    "Standort jederzeit wechseln",
    "KI-Personalisierung an/aus",
    "Sprache: Deutsch / Englisch",
]

txt(slide, "Meine Events", Inches(6.1), Inches(1.9), Inches(6.8), Inches(0.45),
    size=15, bold=True, color=ACC)
bullets(slide, items_my, Inches(6.1), Inches(2.35), Inches(6.7), Inches(1.6),
        size=14, color=GRAY, icon=" •")

txt(slide, "Profil-Einstellungen", Inches(6.1), Inches(4.1), Inches(6.8), Inches(0.45),
    size=15, bold=True, color=ACC)
bullets(slide, items_prof, Inches(6.1), Inches(4.55), Inches(6.7), Inches(1.4),
        size=14, color=GRAY, icon=" •")


# ═══════════════════════════════════════════════
# FOLIE 6 — Organizer-Dashboard
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header(slide, "Organizer-Dashboard", "Events erstellen, bearbeiten, verwalten")

phone_frame(slide, os.path.join(SS, "06_organizer.png"),
            Inches(2.6), H / 2 + Inches(0.3), phone_w=Inches(2.85))

txt(slide, "Fuer Organisatoren",
    Inches(5.55), Inches(1.1), Inches(7.5), Inches(0.6),
    size=26, bold=True, color=DARK)

org_items = [
    ("📋", "Eigene Events", "Alle erstellten Events in einer Liste"),
    ("✏️", "Bearbeiten",    "Titel, Datum, Ort — alles editierbar"),
    ("✕",  "Absagen",       "Event mit einem Klick deaktivieren"),
    ("➕", "Neues Event",   "Formular mit Validierung & Standort-Wahl"),
]
for i, (icon, title, desc) in enumerate(org_items):
    y = Inches(1.85) + i * Inches(1.1)
    rect(slide, Inches(5.55), y, Inches(7.3), Inches(0.95), fill=BG)
    txt(slide, icon, Inches(5.7), y + Inches(0.2), Inches(0.5), Inches(0.55), size=18)
    txt(slide, title, Inches(6.3), y + Inches(0.08), Inches(3.0), Inches(0.42),
        size=15, bold=True, color=DARK)
    txt(slide, desc, Inches(6.3), y + Inches(0.5), Inches(6.3), Inches(0.38),
        size=13, color=GRAY)


# ═══════════════════════════════════════════════
# FOLIE 7 — Abschluss / Call to Action
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, fill=DARK)
rect(slide, 0, H - Inches(0.55), W, Inches(0.55), fill=LIGHT)

# Deko-Kreis
circle = slide.shapes.add_shape(9, int(W - Inches(4.5)), int(-Inches(1.5)),
                                  int(Inches(6)), int(Inches(6)))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x15, 0x50, 0x90)
circle.line.fill.background()

txt(slide, "Bereit zum Ausprobieren?",
    Inches(0.9), Inches(1.5), Inches(8.5), Inches(0.9),
    size=34, bold=True, color=WHITE)

summary = [
    "Feed mit Swipe-Interaktion",
    "Standort- und Typ-Filter",
    "KI-Personalisierung",
    "Organizer- und Admin-Portal",
    "Wartelisten-Management",
]
bullets(slide, summary, Inches(0.9), Inches(2.6), Inches(6.0), Inches(3.0),
        size=18, color=LIGHT, icon=" >")

rect(slide, Inches(0.9), Inches(5.55), Inches(5.2), Inches(0.8), fill=LIGHT)
txt(slide, "  Live-Demo  ->  http://localhost:5173",
    Inches(0.9), Inches(5.6), Inches(5.1), Inches(0.7),
    size=15, bold=True, color=DARK)

txt(slide, "Fragen?",
    Inches(8.5), Inches(3.5), Inches(4.0), Inches(0.8),
    size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, "adesso EventApp  |  Mai 2026",
    Inches(0.9), Inches(6.75), Inches(7), Inches(0.45),
    size=12, color=RGBColor(0x8A, 0xAA, 0xC8))


# ── Speichern ──────────────────────────────────
out = os.path.join(DOCS, "adessoEventApp_Prototype_Demo.pptx")
prs.save(out)
print(f"Gespeichert: {out}")
